from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Entrepreneurial Process and Timmon's Model"
subtitle.text = "Lecture Presentation"

# Slide 2: Learning Objectives
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Lecture Learning Objectives"
content.text = "- Understand the entrepreneurial process\n" \
               "- Identify key factors influencing the decision to start a company\n" \
               "- Identify crucial components for a successful new business\n" \
               "- List different steps involved in the entrepreneurial process"

# Slide 3: The Entrepreneurial Process
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "The Entrepreneurial Process"
content.text = "1. Identifying an opportunity\n" \
               "2. Developing a business concept\n" \
               "3. Acquiring resources\n" \
               "4. Implementing and managing the business\n" \
               "5. Harvesting the venture"

# Slide 4: Key Factors Influencing Entrepreneurship
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Key Factors Influencing Entrepreneurship"
content.text = "- Market demand\n" \
               "- Innovation and creativity\n" \
               "- Access to capital and resources\n" \
               "- Business environment and policies\n" \
               "- Entrepreneurial mindset and risk-taking"

# Slide 5: Components of a Successful New Business
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Components of a Successful New Business"
content.text = "- Strong leadership and vision\n" \
               "- Effective business strategy\n" \
               "- Financial planning and management\n" \
               "- Marketing and customer acquisition\n" \
               "- Adaptability and resilience"

# Slide 6: Timmon’s Model of Entrepreneurship
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Timmon’s Model of Entrepreneurship"
content.text = "Timmon's Model highlights three critical factors:\n\n" \
               "1. **Opportunity** - The foundation of a new venture\n" \
               "2. **Resources** - Financial, human, and social capital\n" \
               "3. **Team** - A strong, committed entrepreneurial team\n\n" \
               "Successful entrepreneurship is balancing these three elements effectively."

# Slide 7: Conclusion
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion"
content.text = "- The entrepreneurial process involves multiple steps from ideation to execution.\n" \
               "- Several factors influence the decision to start a business.\n" \
               "- Timmon’s Model emphasizes opportunity, resources, and team as key elements.\n" \
               "- A well-structured business approach leads to successful entrepreneurship."

# Save the presentation
pptx_file = "E:/python/Entrepreneurial_Process_Timmons_Model.pptx"
prs.save(pptx_file)

# Return the file path for user to download
pptx_file