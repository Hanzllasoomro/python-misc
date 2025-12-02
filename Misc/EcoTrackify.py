from pptx import Presentation  
from pptx.util import Inches  

# Create a new presentation  
prs = Presentation()  

# Slide 1: Title Slide  
slide_layout = prs.slide_layouts[0]  
slide = prs.slides.add_slide(slide_layout)  
title = slide.shapes.title  
subtitle = slide.placeholders[1]  
title.text = "EcoTrackify"  
subtitle.text = "A Personal Carbon Tracking App"  

# Slide 2: Introduction  
slide_layout = prs.slide_layouts[1]  
slide = prs.slides.add_slide(slide_layout)  
title = slide.shapes.title  
content = slide.placeholders[1]  
title.text = "Introduction"  
content.text = "EcoTrackify is designed to help users monitor their carbon footprint by tracking daily activities and providing insights for sustainable living."  

# Slide 3: Features  
slide_layout = prs.slide_layouts[1]  
slide = prs.slides.add_slide(slide_layout)  
title = slide.shapes.title  
content = slide.placeholders[1]  
title.text = "Key Features"  
content.text = "- Daily Carbon Footprint Tracking\n- Activity-Based Emission Calculation\n- User-Friendly Insights & Graphs\n- Personalized Suggestions for Sustainability"  

# Slide 4: Technology Stack  
slide_layout = prs.slide_layouts[1]  
slide = prs.slides.add_slide(slide_layout)  
title = slide.shapes.title  
content = slide.placeholders[1]  
title.text = "Technology Stack"  
content.text = "- Flutter (Frontend)\n- Firebase (Backend & Authentication)\n- REST APIs for Carbon Data\n- AI-Based Recommendation System (Future Enhancement)"  

# Slide 5: Timeline  
slide_layout = prs.slide_layouts[5]  
slide = prs.slides.add_slide(slide_layout)  
title = slide.shapes.title  
title.text = "Project Timeline"  
content = slide.placeholders[1]  
content.text = "• Week 1-2: Research & Planning\n• Week 3-4: UI/UX Design & Prototyping\n• Week 5-7: Backend Development & API Integration\n• Week 8: Testing & Final Deployment"  

# Slide 6: Working Snapshots (Placeholder)  
slide_layout = prs.slide_layouts[5]  
slide = prs.slides.add_slide(slide_layout)  
title = slide.shapes.title  
title.text = "App Screenshots (Placeholder)"  
content = slide.placeholders[1]  
content.text = "Replace this slide with actual screenshots of the app in action."  

# Slide 7: Team Contributions  
slide_layout = prs.slide_layouts[1]  
slide = prs.slides.add_slide(slide_layout)  
title = slide.shapes.title  
content = slide.placeholders[1]  
title.text = "Team Contributions"  
content.text = "- Developer: [Your Name] - Flutter Development, Backend Integration\n- UI/UX Designer: [Teammate Name] - Figma Design & User Experience\n- Data Analyst: [Teammate Name] - Emission Calculations & API Data"  

# Slide 8: Future Enhancements  
slide_layout = prs.slide_layouts[1]  
slide = prs.slides.add_slide(slide_layout)  
title = slide.shapes.title  
content = slide.placeholders[1]  
title.text = "Future Enhancements"  
content.text = "- AI-Powered Emission Predictions\n- Social Features (Compare with Friends)\n- Gamification for Sustainable Habits\n- Integration with Smart Home Devices"  

# Slide 9: Thank You Slide  
slide_layout = prs.slide_layouts[1]  
slide = prs.slides.add_slide(slide_layout)  
title = slide.shapes.title  
content = slide.placeholders[1]  
title.text = "Thank You!"  
content.text = "For more details, contact: [Your Email] | [GitHub Repo]"  

# Save the presentation  
pptx_filename = "./EcoTrackify_Presentation.pptx"  
prs.save(pptx_filename)  

pptx_filename
